# -*- coding: utf-8 -*-
"""生成「班级课表助手」项目答辩 PPT（16:9）。

用法：
    python make_ppt.py
输出：
    ../班级课表助手_答辩PPT.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(BASE, "screenshots")
OUT = os.path.join(os.path.dirname(BASE), "班级课表助手_答辩PPT.pptx")

PRIMARY = RGBColor(0x2B, 0x4E, 0xC9)   # 主色（与产品一致，群青）
DARK = RGBColor(0x1C, 0x1B, 0x1A)      # 墨色标题
TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_LINE = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "PingFang SC"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def _style_run(run, size, color=TEXT, bold=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6):
    """lines: list of (text, size, color, bold) 或 (text, size)。"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, size = item[0], item[1]
        color = item[2] if len(item) > 2 else TEXT
        bold = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        _style_run(run, size, color, bold)
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def title_bar(slide, title, section=""):
    """每页统一标题条：主色横条 + 白色标题。"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.92), fill=PRIMARY)
    add_text(slide, Inches(0.55), Inches(0.0), Inches(9.5), Inches(0.92),
             [(title, 30, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if section:
        add_text(slide, Inches(9.6), Inches(0.0), Inches(3.3), Inches(0.92),
                 [(section, 13, RGBColor(0xE4, 0xEA, 0xF9), False)],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_img(slide, name, x, y, max_w, max_h, caption=None):
    """按原比例把截图放入 (x, y, max_w, max_h) 框内（顶部对齐），加白底+浅灰边框。"""
    path = os.path.join(SHOTS, name)
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    w = max_w
    h = Emu(int(w / ratio))
    if h > max_h:
        h = max_h
        w = Emu(int(h * ratio))
    pad = Inches(0.05)
    add_rect(slide, x - pad, y - pad, w + pad * 2, h + pad * 2,
             fill=WHITE, line=LIGHT_LINE, line_w=1.0)
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    pic.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
    pic.line.width = Pt(0.75)
    if caption:
        add_text(slide, x - pad, y + h + Inches(0.10), w + pad * 2, Inches(1.2),
                 [(caption, 13, TEXT, False)], align=PP_ALIGN.LEFT, space_after=2)
    return w, h


def bullets(slide, x, y, w, h, items, size=16, gap=8):
    """items: (level, text, bold)；level 0 为主点，1 为子点。"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (level, text, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap if level == 0 else gap - 3)
        prefix = "▪ " if level == 0 else "– "
        run = p.add_run()
        run.text = ("    " * level) + prefix + text
        _style_run(run, size if level == 0 else size - 2,
                   DARK if (level == 0 and bold) else TEXT, bold)
    return box


# ============================================================ 1 封面
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=PRIMARY)
add_rect(s, 0, Inches(4.95), SLIDE_W, Inches(0.035), fill=WHITE)
add_text(s, Inches(1.2), Inches(2.1), Inches(10.9), Inches(1.4),
         [("班级课表助手", 48, WHITE, True)])
add_text(s, Inches(1.2), Inches(3.5), Inches(10.9), Inches(0.8),
         [("AI-Coding 实训题目 2 · 项目答辩", 24, WHITE, False)])
add_text(s, Inches(1.2), Inches(5.35), Inches(10.9), Inches(1.2),
         [("汇报小组：王思远 · 许宇坤 · 张耀伊", 18, WHITE, False),
          ("日期：2026 年 8 月", 18, WHITE, False)], space_after=10)

# ============================================ 2 小组成员介绍及分工
s = new_slide()
title_bar(s, "小组成员介绍及分工", "团队")
card_w = Inches(3.95)
gap = Inches(0.25)
x0 = Inches(0.55)
members = [
    ("王思远", "产品负责人", "学号：__________", [
        "主笔《PRD 产品需求文档》v2.0，维护唯一需求上游",
        "拆解 12 个用户故事，定义验收基线",
        "把控需求变更与范围，组织验收",
    ]),
    ("许宇坤", "技术负责人", "学号：__________", [
        "后端任务 B01–B13 共 13 项（认证/课表/导入/权限）",
        "前端任务 F01–F10 共 10 项（课表 SPA/管理后台）",
        "9 表数据模型设计，事务与安全方案落地",
    ]),
    ("张耀伊", "测试负责人", "学号：__________", [
        "测试任务 T01–T10 共 10 项，覆盖接口与回归",
        "70 项 pytest 自动化测试全部通过",
        "事务回滚测试、缺陷跟踪与验收核对",
    ]),
]
for i, (name, role, sid, duties) in enumerate(members):
    x = x0 + i * (card_w + gap)
    add_rect(s, x, Inches(1.35), card_w, Inches(5.35), fill=WHITE, line=LIGHT_LINE, line_w=1.0)
    add_rect(s, x, Inches(1.35), card_w, Inches(0.14), fill=PRIMARY)
    add_text(s, x + Inches(0.3), Inches(1.75), card_w - Inches(0.6), Inches(1.5),
             [(name, 24, DARK, True), (role, 15, PRIMARY, True), (sid, 13, GRAY, False)],
             space_after=4)
    bullets(s, x + Inches(0.3), Inches(3.35), card_w - Inches(0.6), Inches(3.2),
            [(0, d, False) for d in duties], size=13, gap=7)

# ============================================ 3 项目背景与开发思路
s = new_slide()
title_bar(s, "项目背景与开发思路", "背景")
bullets(s, Inches(0.55), Inches(1.25), Inches(6.1), Inches(5.9), [
    (0, "痛点", True),
    (1, "课表分散：截图、Excel、群文件各自为战，版本不一致", False),
    (1, "作业靠问：截止日期口口相传，容易遗漏", False),
    (1, "教务系统面向个人，不支持班级共享课表与作业", False),
    (0, "目标", True),
    (1, "一班一表：全班共享同一份权威课表与作业清单", False),
    (1, "课表即提醒：考试倒计时、作业临期三色提醒", False),
], size=16, gap=10)
bullets(s, Inches(6.95), Inches(1.25), Inches(5.85), Inches(5.9), [
    (0, "技术栈：Flask 3.1 + SQLite + 原生 HTML/CSS/JS 单页前端", True),
    (0, "开发思路", True),
    (1, "PRD 先行 → 契约冻结 → 前后端并行 → 集成联调 → 测试验收", False),
    (0, "架构要点", True),
    (1, "9 表数据模型；课程与上课安排分表，支持一课多时段", False),
    (1, "三角色权限（系统管理员 / 班级管理员 / 学生）", False),
    (1, "服务端为唯一权威：校验、冲突检测、审计均在服务端", False),
], size=16, gap=10)

# ============================================ 4 在线演示
s = new_slide()
title_bar(s, "在线演示", "成果展示")
rows = [
    ("角色", "账号", "密码", "入口 / 说明"),
    ("系统管理员", "sysadmin", "sysadmin123", "系统管理后台 /admin"),
    ("班级管理员", "admin2301 / admin2302", "admin1234", "软件技术 2301 / 2302 班"),
    ("学生（已激活）", "2023301001–005 / 2023302001–005", "student123", "学生周视图 / 日视图"),
    ("学生（待激活）", "2023301006–008 等", "激活时自行设置", "凭班级邀请码 + 姓名 + 学号激活"),
]
tbl_x, tbl_y = Inches(0.55), Inches(1.3)
tbl_w = Inches(12.25)
tbl = s.shapes.add_table(len(rows), 4, tbl_x, tbl_y, tbl_w, Inches(3.1)).table
tbl.columns[0].width = Inches(2.3)
tbl.columns[1].width = Inches(3.7)
tbl.columns[2].width = Inches(2.3)
tbl.columns[3].width = Inches(3.95)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = val
        _style_run(run, 14 if r else 15, WHITE if r == 0 else TEXT, r == 0)
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
add_text(s, Inches(0.55), Inches(4.75), Inches(12.25), Inches(2.3), [
    ("演示路径：sysadmin 看后台 → admin2301 维护课表/导入 CSV/发邀请码 → 学生 2023301001 查看周视图、倒计时与作业提醒", 15, TEXT, False),
    ("启动方式：python app.py 后访问 http://127.0.0.1:5000/ ，管理后台为 /admin", 15, TEXT, False),
    ("⚠ 以上账号与邀请码仅限开发环境演示，禁止用于生产；数据库只保存邀请码 sha256 摘要", 14, RGBColor(0xC0, 0x39, 0x2B), True),
], space_after=10)

# ============================================ 5 登录与学生激活
s = new_slide()
title_bar(s, "成果展示 ①：登录与学生激活", "成果展示")
add_img(s, "01_login.png", Inches(0.7), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="统一登录页：三角色同一入口，按角色跳转课表端或管理后台")
add_img(s, "02_activate.png", Inches(6.95), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="学生激活：邀请码 + 姓名 + 学号校验，自设密码后即为正式账号")

# ===================================== 6 学生周视图 / 日视图
s = new_slide()
title_bar(s, "成果展示 ②：学生课表周视图 / 日视图", "成果展示")
add_img(s, "03_week_view.png", Inches(0.7), Inches(1.35), Inches(6.0), Inches(4.6),
        caption="周视图：考试倒计时 + 作业三色提醒 + 彩色课表，当天列高亮")
add_img(s, "04_day_view.png", Inches(7.6), Inches(1.35), Inches(5.2), Inches(3.3),
        caption="日视图：按节次列出当天课程与作业")

# ===================================== 7 课程详情只读 + PDF 导出
s = new_slide()
title_bar(s, "成果展示 ③：课程详情与 PDF 打印导出", "成果展示")
add_img(s, "05_course_detail_student.png", Inches(0.7), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="学生端课程详情：教师、教室、周次、考试日期等只读展示")
add_img(s, "15_print_view.png", Inches(6.95), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="一键打印横向课表：专用打印样式，可直接导出 PDF 张贴")

# ===================================== 8 课程编辑 + CSV 导入
s = new_slide()
title_bar(s, "成果展示 ④：班级管理员 · 课程维护与导入", "成果展示")
add_img(s, "07_course_edit.png", Inches(0.7), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="课程编辑：一课多时段（如高等数学周一 1-2 节 + 周三 3-4 节）")
add_img(s, "08_import_csv.png", Inches(6.95), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="双 CSV 导入：课程 + 作业统一校验，失败整体回滚，确认后同事务覆盖")

# ===================================== 9 作业管理 + 邀请码
s = new_slide()
title_bar(s, "成果展示 ⑤：班级管理员 · 作业与邀请码", "成果展示")
add_img(s, "10_homework.png", Inches(0.7), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="作业管理：按课程登记作业与截止日期，学生端自动三色提醒")
add_img(s, "09_invite_code.png", Inches(6.95), Inches(1.35), Inches(5.9), Inches(3.7),
        caption="动态邀请码：可重新生成，明文仅一次性展示，库中只存摘要")

# ===================================== 10 系统管理后台
s = new_slide()
title_bar(s, "成果展示 ⑥：系统管理后台", "成果展示")
add_img(s, "11_admin_semesters.png", Inches(0.6), Inches(1.4), Inches(4.1), Inches(2.6),
        caption="学期管理：唯一 active 学期，周数驱动周视图")
add_img(s, "13_admin_users.png", Inches(4.95), Inches(1.4), Inches(3.4), Inches(2.6),
        caption="学生名单 / 用户管理：CSV 导入名单，待激活状态")
add_img(s, "14_admin_audit.png", Inches(8.6), Inches(1.4), Inches(4.1), Inches(2.6),
        caption="审计日志：登录、导入、邀请码等关键操作留痕")

# ===================================== 11 移动端响应式
s = new_slide()
title_bar(s, "成果展示 ⑦：移动端响应式", "成果展示")
add_img(s, "16_mobile.png", Inches(1.1), Inches(1.35), Inches(3.2), Inches(5.5))
bullets(s, Inches(5.2), Inches(1.8), Inches(7.5), Inches(4.5), [
    (0, "同一套页面自适应手机窄屏，无需独立 App", True),
    (0, "课表、倒计时、作业提醒在移动端完整可用", False),
    (0, "学生掏出手机即可查看本周课程与临期作业", False),
    (0, "配合打印样式，覆盖「手机随时查 + 纸质张贴看」两类场景", False),
], size=17, gap=14)

# ===================================== 12 亮点：功能创新
s = new_slide()
title_bar(s, "项目亮点与创新 ①：功能创新", "亮点")
bullets(s, Inches(0.55), Inches(1.4), Inches(12.25), Inches(5.5), [
    (0, "课表即提醒：考试倒计时 + 作业按临期程度三色提醒，课表不只是“看”的", True),
    (0, "双 CSV 原子覆盖导入：课程与作业两份文件统一校验，任一失败整体回滚，确认后在同一事务中覆盖本班当前学期数据", True),
    (0, "一课多时段：课程与上课安排分表，一门课可挂多个星期/节次/教室，贴近真实排课", True),
    (0, "动态邀请码：可重新生成、可停用；数据库只存 sha256 摘要，明文仅生成时一次性展示", True),
], size=18, gap=16)

# ===================================== 13 亮点：技术方案
s = new_slide()
title_bar(s, "项目亮点与创新 ②：技术方案", "亮点")
bullets(s, Inches(0.55), Inches(1.4), Inches(12.25), Inches(5.5), [
    (0, "服务端冲突检测：排课/调整时检测时间冲突，返回 409 与冲突详情，而非静默覆盖", True),
    (0, "条件唯一索引兜底业务规则：每班至多一名管理员、全校唯一 active 学期，由数据库层保证", True),
    (0, "安全基线：所有写请求强制 CSRF 头校验、登录限流、关键操作审计日志", True),
    (0, "注入与 XSS 防护：全程参数化查询，前端输出统一转义", True),
], size=18, gap=16)

# ===================================== 14 亮点：工程质量
s = new_slide()
title_bar(s, "项目亮点与创新 ③：工程质量", "亮点")
bullets(s, Inches(0.55), Inches(1.4), Inches(12.25), Inches(5.5), [
    (0, "70 项 pytest 自动化测试全部通过，覆盖认证、权限、导入、冲突等核心路径", True),
    (0, "事务与回滚测试：专门验证导入失败时数据库不产生半成品数据", True),
    (0, "幂等建表与种子脚本：schema 可重复执行，演示数据一键重灌", True),
    (0, "响应式布局 + 独立打印样式：手机端与纸质课表同源输出", True),
], size=18, gap=16)

# ===================================== 15 学习方法分享
s = new_slide()
title_bar(s, "学习方法分享", "方法")
bullets(s, Inches(0.55), Inches(1.3), Inches(12.25), Inches(3.3), [
    (0, "AI 辅助编程方法", True),
    (1, "PRD 作为唯一需求上游：先让 AI 参与打磨需求文档，再谈写代码", False),
    (1, "先冻结接口契约，前后端再并行开发，联调按契约对表", False),
    (1, "用自动化测试验证 AI 产出的代码，而不是目测“看起来对”", False),
], size=16, gap=8)
bullets(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(2.8), [
    (0, "问题解决举例：本机 LibreSSL 不提供 scrypt，密码哈希直接崩溃 → 显式指定 pbkdf2 并固定参数，问题定位到环境差异后一行配置解决", True),
    (0, "工具与资源：Flask / SQLite 官方文档、pytest、Playwright、python-pptx", True),
], size=16, gap=8)

# ===================================== 16 实习收获与感受
s = new_slide()
title_bar(s, "实习收获与感受分享", "收获")
add_text(s, Inches(0.55), Inches(1.15), Inches(12.25), Inches(0.5),
         [("（示例文案，可按真实感受修改）", 13, GRAY, False)])
bullets(s, Inches(0.55), Inches(1.75), Inches(12.25), Inches(5.2), [
    (0, "专业技能：体验了一次完整的全栈契约化开发——从 PRD、接口契约到事务实现，理解了“数据一致性靠事务与约束兜底，而不是靠小心”", True),
    (0, "软实力：文档先行让分工协作有了共同语言，跨角色沟通（产品/技术/测试）围绕验收基线展开，返工明显减少", True),
    (0, "认知变化：AI 是放大器——能把清晰的需求快速变成代码，但验收标准必须由人来定，测试写不出来往往说明需求没想清", True),
], size=17, gap=16)

# ===================================== 17 结尾
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=PRIMARY)
add_text(s, Inches(1.2), Inches(2.6), Inches(10.9), Inches(1.3),
         [("谢谢聆听", 44, WHITE, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(4.0), Inches(10.9), Inches(0.9),
         [("Q & A  欢迎提问与交流", 24, WHITE, False)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(5.3), Inches(10.9), Inches(0.6),
         [("班级课表助手 · AI-Coding 实训题目 2", 16, WHITE, False)], align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"saved: {OUT}, slides: {len(prs.slides._sldIdLst)}")
